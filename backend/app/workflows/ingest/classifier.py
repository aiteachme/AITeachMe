"""Lightweight ingest classification used to plan parser routing."""

from __future__ import annotations

from pathlib import Path
import re
from typing import Any

from pydantic import BaseModel, Field
import structlog


try:
    import fitz
except ImportError:
    fitz = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


logger = structlog.get_logger()

_ZH_RE = re.compile(r"[\u4e00-\u9fff]")
_EN_RE = re.compile(r"[a-zA-Z]")
_HEADING_LIKE_RE = re.compile(
    r"^(第[一二三四五六七八九十百千万\d]+[章节篇]|Chapter\s+\d+|Section\s+\d+|\d+[\.\s])",
    re.MULTILINE,
)
_TABLE_LIKE_RE = re.compile(r"\|.*\|.*\||\+[-=]+\+", re.MULTILINE)
_FORMULA_RE = re.compile(r"[≈≠≤≥∑∫∞√πΔλμσθ]|\\frac|\\sum|\\int")
_IMAGE_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff"})


class ClassificationResult(BaseModel):
    """Classification output consumed by the ingest workflow."""

    file_category: str
    text_density: float = 0.0
    ocr_ratio: float = 0.0
    image_page_ratio: float = 0.0
    heading_count: int = 0
    estimated_pages: int = 0
    detected_language: str = "unknown"
    has_tables: bool = False
    has_formulas: bool = False
    recommended_parser: str = ""
    fallback_parsers: list[str] = Field(default_factory=list)

    def to_dict(self) -> dict[str, object]:
        return self.model_dump()


def classify_file(file_path: str | Path, filetype: str) -> ClassificationResult:
    """Inspect a file quickly and choose a parser chain."""

    path = Path(file_path)
    extension = filetype.lower() if filetype.startswith(".") else f".{filetype.lower()}"

    if extension == ".pdf":
        return _classify_pdf(path)
    if extension in {".ppt", ".pptx"}:
        return _classify_pptx(path)
    if extension == ".docx":
        return _classify_docx(path)
    if extension in _IMAGE_EXTENSIONS:
        return ClassificationResult(
            file_category="image",
            estimated_pages=1,
            recommended_parser="llm_vision",
        )
    return ClassificationResult(
        file_category="unknown",
        recommended_parser="markitdown",
    )


def _classify_pdf(path: Path) -> ClassificationResult:
    if fitz is None:
        return ClassificationResult(
            file_category="text_pdf",
            recommended_parser="pymupdf4llm",
            fallback_parsers=["markitdown", "pymupdf_native"],
        )

    document = fitz.open(str(path))
    total_pages = len(document)
    sample_pages = min(total_pages, 10)

    total_chars = 0
    image_heavy_pages = 0
    all_text_parts: list[str] = []

    for page_index in range(sample_pages):
        page = document[page_index]
        text = page.get_text("text") or ""
        char_count = len(text.strip())
        total_chars += char_count

        images = page.get_images(full=True)
        large_images = [image for image in images if _is_large_image(document, image[0])]
        if char_count < 50 and large_images:
            image_heavy_pages += 1

        all_text_parts.append(text)

    document.close()

    all_text = "".join(all_text_parts)
    avg_density = total_chars / sample_pages if sample_pages else 0
    image_ratio = image_heavy_pages / sample_pages if sample_pages else 0
    detected_language = _detect_language(all_text[:5000])
    heading_count = len(_HEADING_LIKE_RE.findall(all_text[:10000]))
    has_tables = bool(_TABLE_LIKE_RE.search(all_text[:10000]))
    has_formulas = bool(_FORMULA_RE.search(all_text[:10000]))

    if avg_density < 30:
        file_category = "scanned_pdf"
        recommended_parser = "pymupdf_native"
        fallback_parsers = ["markitdown"]
    elif image_ratio > 0.5:
        file_category = "complex_pdf"
        recommended_parser = "pymupdf4llm"
        fallback_parsers = ["pymupdf_native", "markitdown"]
    else:
        file_category = "text_pdf"
        recommended_parser = "pymupdf4llm"
        fallback_parsers = ["markitdown", "pymupdf_native"]

    result = ClassificationResult(
        file_category=file_category,
        text_density=round(avg_density, 1),
        ocr_ratio=round(image_ratio, 2),
        image_page_ratio=round(image_ratio, 2),
        heading_count=heading_count,
        estimated_pages=total_pages,
        detected_language=detected_language,
        has_tables=has_tables,
        has_formulas=has_formulas,
        recommended_parser=recommended_parser,
        fallback_parsers=fallback_parsers,
    )
    logger.info(
        "classify_pdf_done",
        filename=path.name,
        category=result.file_category,
        pages=result.estimated_pages,
        avg_density=result.text_density,
        image_ratio=result.image_page_ratio,
        language=result.detected_language,
    )
    return result


def _is_large_image(document: Any, xref: int) -> bool:
    try:
        image = document.extract_image(xref)
    except Exception:
        return False
    if image is None:
        return False
    return len(image["image"]) > 5000


def _classify_pptx(path: Path) -> ClassificationResult:
    if Presentation is None:
        return ClassificationResult(
            file_category="pptx",
            recommended_parser="markitdown",
            fallback_parsers=["python_pptx_native"],
        )

    try:
        presentation = Presentation(str(path))
        slide_count = len(presentation.slides)
        total_text = 0
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                total_text += sum(len(paragraph.text.strip()) for paragraph in shape.text_frame.paragraphs)
        avg_density = total_text / slide_count if slide_count else 0
    except Exception:
        slide_count = 0
        avg_density = 0

    return ClassificationResult(
        file_category="pptx",
        text_density=round(avg_density, 1),
        estimated_pages=slide_count,
        recommended_parser="markitdown",
        fallback_parsers=["python_pptx_native"],
    )


def _classify_docx(path: Path) -> ClassificationResult:
    if Document is None:
        return ClassificationResult(
            file_category="docx",
            recommended_parser="markitdown",
            fallback_parsers=["python_docx_native"],
        )

    try:
        document = Document(str(path))
        paragraph_count = len(document.paragraphs)
        total_text = sum(len(paragraph.text.strip()) for paragraph in document.paragraphs)
        heading_count = sum(
            1
            for paragraph in document.paragraphs
            if paragraph.style and paragraph.style.name and paragraph.style.name.startswith("Heading")
        )
        avg_density = total_text / max(paragraph_count, 1)
    except Exception:
        paragraph_count = 0
        heading_count = 0
        avg_density = 0

    return ClassificationResult(
        file_category="docx",
        text_density=round(avg_density, 1),
        heading_count=heading_count,
        estimated_pages=max(paragraph_count // 30, 1),
        recommended_parser="markitdown",
        fallback_parsers=["python_docx_native"],
    )


def _detect_language(text: str) -> str:
    zh_count = len(_ZH_RE.findall(text))
    en_count = len(_EN_RE.findall(text))
    if zh_count > en_count * 2:
        return "zh"
    if en_count > zh_count * 2:
        return "en"
    return "mixed"
