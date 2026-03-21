"""PPTX parser variants used by the ingest workflow."""

from __future__ import annotations

from collections.abc import Iterable, Iterator
from pathlib import Path

import structlog

from app.core.exceptions import FileParseError
from app.workflows.ingest.parse_utils import save_image_bytes


try:
    from markitdown import MarkItDown
except ImportError:
    MarkItDown = None

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None
    MSO_SHAPE_TYPE = None


logger = structlog.get_logger()


async def parse_pptx_with_markitdown(file_path: str | Path, asset_dir: Path) -> str:
    """Parse PPTX with MarkItDown and supplement image extraction."""

    path = Path(file_path)
    if MarkItDown is None:
        raise FileParseError(path.name, reason="MarkItDown is not available.")

    logger.info("parse_pptx_markitdown_start", filename=path.name)
    result = MarkItDown().convert(str(path))
    if not result.text_content or not result.text_content.strip():
        raise FileParseError(path.name, reason="MarkItDown returned empty markdown.")

    supplement_pptx_images(path, asset_dir)
    return result.text_content


async def parse_pptx_with_python_pptx(file_path: str | Path, asset_dir: Path) -> str:
    """Parse PPTX with python-pptx and inline extracted images."""

    path = Path(file_path)
    if Presentation is None or MSO_SHAPE_TYPE is None:
        raise FileParseError(path.name, reason="python-pptx is not available.")

    logger.info("parse_pptx_python_native_start", filename=path.name)
    presentation = Presentation(str(path))
    sections: list[str] = []
    image_count = 0

    for slide_index, slide in enumerate(presentation.slides, start=1):
        sections.append(f"## Slide {slide_index}")
        for shape in _iter_shapes(slide.shapes):
            if shape.has_text_frame:
                sections.extend(
                    paragraph.text.strip()
                    for paragraph in shape.text_frame.paragraphs
                    if paragraph.text.strip()
                )

            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            image_bytes = shape.image.blob
            if len(image_bytes) < 1024:
                continue

            image_count += 1
            ext = _normalize_image_extension(shape.image.content_type)
            filename = save_image_bytes(
                image_bytes,
                asset_dir,
                name_hint=f"slide{slide_index}_img{image_count}",
                ext=f".{ext}",
            )
            sections.append(f"![Slide {slide_index} image {image_count}]({filename})")

        sections.append("")

    result = "\n".join(sections).strip()
    if not result:
        raise FileParseError(path.name, reason="PPTX text extraction returned empty markdown.")

    logger.info("parse_pptx_python_native_done", filename=path.name, images=image_count)
    return result


def supplement_pptx_images(file_path: Path, asset_dir: Path) -> None:
    """Extract PPTX images when markdown came from MarkItDown."""

    if Presentation is None or MSO_SHAPE_TYPE is None:
        return

    presentation = Presentation(str(file_path))
    count = 0
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in _iter_shapes(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            image_bytes = shape.image.blob
            if len(image_bytes) < 1024:
                continue

            save_image_bytes(
                image_bytes,
                asset_dir,
                name_hint=f"slide{slide_index}_img{count + 1}",
                ext=f".{_normalize_image_extension(shape.image.content_type)}",
            )
            count += 1

    if count:
        logger.info("pptx_images_supplemented", filename=file_path.name, count=count)


def _iter_shapes(shapes: Iterable[object]) -> Iterator[object]:
    for shape in shapes:
        yield shape
        if MSO_SHAPE_TYPE is None or shape.shape_type != MSO_SHAPE_TYPE.GROUP:
            continue
        yield from _iter_shapes(shape.shapes)


def _normalize_image_extension(content_type: str | None) -> str:
    if not content_type:
        return "png"
    extension = content_type.split("/")[-1]
    return "jpg" if extension == "jpeg" else extension
