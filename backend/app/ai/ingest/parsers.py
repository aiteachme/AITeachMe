"""
文档解析器 — PDF / PPTX / DOCX / Image → Markdown

所有解析器输出 CommonMark 兼容的 Markdown。
图片解析通过多模态 LLM 调用，失败时由调用方标记 parse_failed。

需求：5.4, 5.5, 5.6, 5.7, 6.1
"""

from __future__ import annotations

from pathlib import Path

import structlog

from app.core.exceptions import FileParseError

logger = structlog.get_logger()


async def parse_pdf(file_path: str | Path) -> str:
    """PDF → Markdown，使用 MarkItDown（优先）回退到 PyMuPDF4LLM。"""
    file_path = Path(file_path)
    logger.info("parse_pdf_start", file=file_path.name)

    # 优先使用 markitdown
    try:
        from markitdown import MarkItDown

        md = MarkItDown()
        result = md.convert(str(file_path))
        text = result.text_content
        if text and text.strip():
            logger.info("parse_pdf_done", file=file_path.name, method="markitdown", length=len(text))
            return text
    except Exception as exc:
        logger.warning("parse_pdf_markitdown_failed", file=file_path.name, error=str(exc))

    # 回退到 pymupdf4llm
    try:
        import pymupdf4llm

        text = pymupdf4llm.to_markdown(str(file_path))
        if text and text.strip():
            logger.info("parse_pdf_done", file=file_path.name, method="pymupdf4llm", length=len(text))
            return text
    except Exception as exc:
        logger.warning("parse_pdf_pymupdf4llm_failed", file=file_path.name, error=str(exc))

    raise FileParseError(file_path.name, reason="PDF 解析失败：markitdown 和 pymupdf4llm 均无法提取内容")


async def parse_pptx(file_path: str | Path) -> str:
    """PPT/PPTX → Markdown，提取幻灯片文本 + 演讲者备注，按标题层级结构化。"""
    file_path = Path(file_path)
    logger.info("parse_pptx_start", file=file_path.name)

    try:
        from markitdown import MarkItDown

        md = MarkItDown()
        result = md.convert(str(file_path))
        text = result.text_content
        if text and text.strip():
            logger.info("parse_pptx_done", file=file_path.name, length=len(text))
            return text
    except Exception as exc:
        logger.warning("parse_pptx_markitdown_failed", file=file_path.name, error=str(exc))

    # 回退：使用 python-pptx 手动提取
    try:
        from pptx import Presentation

        prs = Presentation(str(file_path))
        parts: list[str] = []
        for idx, slide in enumerate(prs.slides, 1):
            parts.append(f"## 幻灯片 {idx}")
            # 提取文本框内容
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            parts.append(line)
            # 提取演讲者备注
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append("")
                    parts.append(f"> **备注：** {notes}")
            parts.append("")

        text = "\n".join(parts)
        if text.strip():
            logger.info("parse_pptx_done", file=file_path.name, method="python-pptx", length=len(text))
            return text
    except Exception as exc:
        logger.warning("parse_pptx_fallback_failed", file=file_path.name, error=str(exc))

    raise FileParseError(file_path.name, reason="PPTX 解析失败")


async def parse_docx(file_path: str | Path) -> str:
    """DOCX → Markdown，使用 MarkItDown。"""
    file_path = Path(file_path)
    logger.info("parse_docx_start", file=file_path.name)

    try:
        from markitdown import MarkItDown

        md = MarkItDown()
        result = md.convert(str(file_path))
        text = result.text_content
        if text and text.strip():
            logger.info("parse_docx_done", file=file_path.name, length=len(text))
            return text
    except Exception as exc:
        logger.warning("parse_docx_failed", file=file_path.name, error=str(exc))

    raise FileParseError(file_path.name, reason="DOCX 解析失败")


async def parse_image(file_path: str | Path) -> str:
    """图片（PNG/JPG/WEBP）→ Markdown，通过多模态 LLM 描述图片内容。

    LLM 调用失败时抛出 FileParseError，由调用方标记 parse_failed。
    """
    import base64

    file_path = Path(file_path)
    logger.info("parse_image_start", file=file_path.name)

    suffix = file_path.suffix.lower()
    mime_map = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".webp": "image/webp"}
    mime_type = mime_map.get(suffix, "image/png")

    # 读取并 base64 编码
    image_data = file_path.read_bytes()
    b64 = base64.b64encode(image_data).decode("utf-8")

    messages = [
        {
            "role": "user",
            "content": [
                {
                    "type": "text",
                    "text": (
                        "请详细描述这张图片的内容，将其转化为结构化的 Markdown 文本。"
                        "如果图片包含文字，请完整提取。如果是图表，请描述其结构和数据。"
                    ),
                },
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:{mime_type};base64,{b64}"},
                },
            ],
        }
    ]

    try:
        from app.core.llm import acompletion

        text = await acompletion(messages)
        if text and text.strip():
            logger.info("parse_image_done", file=file_path.name, length=len(text))
            return text
    except Exception as exc:
        logger.error("parse_image_llm_failed", file=file_path.name, error=str(exc))
        raise FileParseError(file_path.name, reason=f"图片 LLM 解析失败：{exc}") from exc

    raise FileParseError(file_path.name, reason="图片 LLM 返回空内容")
