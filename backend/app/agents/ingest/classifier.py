"""文件轻量分类器：纯规则 + 统计，不调用 LLM。

根据文件类型和内容特征，判断文件类别、推荐 parser 和 fallback 链。
"""

from __future__ import annotations

import re
from dataclasses import asdict, dataclass, field
from pathlib import Path

import structlog

logger = structlog.get_logger()


@dataclass
class ClassificationResult:
    """分类结果。"""

    file_category: str  # text_pdf, scanned_pdf, complex_pdf, pptx, docx, image
    text_density: float = 0.0  # 每页平均字符数
    ocr_ratio: float = 0.0  # 估计 OCR 内容占比
    image_page_ratio: float = 0.0  # 图片页占比
    heading_count: int = 0  # 检测到的标题数
    estimated_pages: int = 0  # 页数 / slide 数
    detected_language: str = "unknown"  # zh / en / mixed
    has_tables: bool = False
    has_formulas: bool = False
    recommended_parser: str = ""  # 推荐的主 parser
    fallback_parsers: list[str] = field(default_factory=list)

    def to_dict(self) -> dict:
        return asdict(self)


def classify_file(file_path: str | Path, filetype: str) -> ClassificationResult:
    """对文件做轻量分类，返回分类结果。"""
    path = Path(file_path)
    ext = filetype.lower() if filetype.startswith(".") else f".{filetype.lower()}"

    if ext == ".pdf":
        return _classify_pdf(path)
    elif ext in (".ppt", ".pptx"):
        return _classify_pptx(path)
    elif ext == ".docx":
        return _classify_docx(path)
    elif ext in (".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"):
        return ClassificationResult(
            file_category="image",
            estimated_pages=1,
            recommended_parser="llm_vision",
            fallback_parsers=[],
        )
    else:
        return ClassificationResult(
            file_category="unknown",
            recommended_parser="markitdown",
            fallback_parsers=[],
        )


# ---------------------------------------------------------------------------
# PDF 分类
# ---------------------------------------------------------------------------

_ZH_RE = re.compile(r"[\u4e00-\u9fff]")
_EN_RE = re.compile(r"[a-zA-Z]")
_HEADING_LIKE_RE = re.compile(r"^(第[一二三四五六七八九十百千\d]+[章节篇]|Chapter\s+\d+|Section\s+\d+|\d+[\.\s])", re.MULTILINE)
_TABLE_LIKE_RE = re.compile(r"\|.*\|.*\||\+[-=]+\+", re.MULTILINE)
_FORMULA_RE = re.compile(r"[∑∫∂∇≈≠≤≥±×÷√∞αβγδεζηθλμπσφψω]|\\frac|\\sum|\\int")


def _classify_pdf(path: Path) -> ClassificationResult:
    """PDF 分类：采样前几页判断文本密度、是否扫描件等。"""
    try:
        import fitz
    except ImportError:
        return ClassificationResult(
            file_category="text_pdf",
            recommended_parser="pymupdf4llm",
            fallback_parsers=["markitdown", "pymupdf_native"],
        )

    doc = fitz.open(str(path))
    total_pages = len(doc)
    sample_pages = min(total_pages, 10)

    total_chars = 0
    image_heavy_pages = 0
    all_text = ""

    for i in range(sample_pages):
        page = doc[i]
        text = page.get_text("text") or ""
        char_count = len(text.strip())
        total_chars += char_count

        # 判断是否图片为主的页面
        images = page.get_images(full=True)
        large_images = [img for img in images if _is_large_image(doc, img[0])]
        if char_count < 50 and len(large_images) > 0:
            image_heavy_pages += 1

        all_text += text

    doc.close()

    avg_density = total_chars / sample_pages if sample_pages > 0 else 0
    image_ratio = image_heavy_pages / sample_pages if sample_pages > 0 else 0

    # 语言检测
    zh_count = len(_ZH_RE.findall(all_text[:5000]))
    en_count = len(_EN_RE.findall(all_text[:5000]))
    if zh_count > en_count * 2:
        lang = "zh"
    elif en_count > zh_count * 2:
        lang = "en"
    else:
        lang = "mixed"

    heading_count = len(_HEADING_LIKE_RE.findall(all_text[:10000]))
    has_tables = bool(_TABLE_LIKE_RE.search(all_text[:10000]))
    has_formulas = bool(_FORMULA_RE.search(all_text[:10000]))

    # 分类决策
    if avg_density < 30:
        category = "scanned_pdf"
        recommended = "pymupdf_native"
        fallbacks = ["markitdown"]
    elif image_ratio > 0.5:
        category = "complex_pdf"
        recommended = "pymupdf4llm"
        fallbacks = ["pymupdf_native", "markitdown"]
    else:
        category = "text_pdf"
        recommended = "pymupdf4llm"
        fallbacks = ["markitdown", "pymupdf_native"]

    result = ClassificationResult(
        file_category=category,
        text_density=round(avg_density, 1),
        ocr_ratio=round(image_ratio, 2),
        image_page_ratio=round(image_ratio, 2),
        heading_count=heading_count,
        estimated_pages=total_pages,
        detected_language=lang,
        has_tables=has_tables,
        has_formulas=has_formulas,
        recommended_parser=recommended,
        fallback_parsers=fallbacks,
    )

    logger.info(
        "classify_pdf_done",
        filename=path.name,
        category=category,
        pages=total_pages,
        avg_density=round(avg_density, 1),
        image_ratio=round(image_ratio, 2),
        language=lang,
    )
    return result


def _is_large_image(doc: object, xref: int) -> bool:
    """判断图片是否足够大（非装饰性）。"""
    try:
        img = doc.extract_image(xref)  # type: ignore[union-attr]
        if img is None:
            return False
        return len(img["image"]) > 5000
    except Exception:
        return False


# ---------------------------------------------------------------------------
# PPTX 分类
# ---------------------------------------------------------------------------


def _classify_pptx(path: Path) -> ClassificationResult:
    """PPTX 分类。"""
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        slide_count = len(prs.slides)
        total_text = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        total_text += len(para.text.strip())

        avg_density = total_text / slide_count if slide_count > 0 else 0
    except Exception:
        slide_count = 0
        avg_density = 0

    return ClassificationResult(
        file_category="pptx",
        text_density=round(avg_density, 1),
        estimated_pages=slide_count,
        recommended_parser="markitdown",
        fallback_parsers=["python_pptx_native"],
    )


# ---------------------------------------------------------------------------
# DOCX 分类
# ---------------------------------------------------------------------------


def _classify_docx(path: Path) -> ClassificationResult:
    """DOCX 分类。"""
    try:
        from docx import Document

        doc = Document(str(path))
        para_count = len(doc.paragraphs)
        total_text = sum(len(p.text.strip()) for p in doc.paragraphs)
        heading_count = sum(
            1 for p in doc.paragraphs
            if p.style and p.style.name and p.style.name.startswith("Heading")
        )
        avg_density = total_text / max(para_count, 1)
    except Exception:
        para_count = 0
        total_text = 0
        heading_count = 0
        avg_density = 0

    return ClassificationResult(
        file_category="docx",
        text_density=round(avg_density, 1),
        heading_count=heading_count,
        estimated_pages=max(para_count // 30, 1),
        recommended_parser="markitdown",
        fallback_parsers=["python_docx_native"],
    )
