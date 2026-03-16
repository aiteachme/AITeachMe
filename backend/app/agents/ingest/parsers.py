"""文件解析器集合。"""

from __future__ import annotations

import base64
from pathlib import Path

import structlog

from app.agents.ingest.prompts import SYSTEM_PROMPT_IMAGE_PARSE
from app.core.exceptions import FileParseError
from app.core.prompt_loader import populate_prompt
from app.schemas.llm import ChatMessage, USER

logger = structlog.get_logger()


async def parse_pdf(file_path: str | Path) -> str:
    """把 PDF 解析为 Markdown。"""

    path = Path(file_path)
    logger.info("parse_pdf_start", filename=path.name)

    try:
        from markitdown import MarkItDown

        result = MarkItDown().convert(str(path))
        if result.text_content and result.text_content.strip():
            return result.text_content
    except Exception as exc:
        logger.warning("parse_pdf_markitdown_failed", filename=path.name, error=str(exc))

    try:
        import pymupdf4llm

        text = pymupdf4llm.to_markdown(str(path))
        if text and text.strip():
            return text
    except Exception as exc:
        logger.warning("parse_pdf_pymupdf4llm_failed", filename=path.name, error=str(exc))

    raise FileParseError(path.name, reason="PDF 解析失败。")


async def parse_pptx(file_path: str | Path) -> str:
    """把 PPT/PPTX 解析为 Markdown。"""

    path = Path(file_path)
    logger.info("parse_pptx_start", filename=path.name)

    try:
        from markitdown import MarkItDown

        result = MarkItDown().convert(str(path))
        if result.text_content and result.text_content.strip():
            return result.text_content
    except Exception as exc:
        logger.warning("parse_pptx_markitdown_failed", filename=path.name, error=str(exc))

    try:
        from pptx import Presentation

        presentation = Presentation(str(path))
        sections: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            sections.append(f"## 幻灯片 {index}")
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        sections.append(text)
            sections.append("")
        return "\n".join(sections).strip()
    except Exception as exc:
        logger.warning("parse_pptx_fallback_failed", filename=path.name, error=str(exc))

    raise FileParseError(path.name, reason="PPTX 解析失败。")


async def parse_docx(file_path: str | Path) -> str:
    """把 DOCX 解析为 Markdown。"""

    path = Path(file_path)
    logger.info("parse_docx_start", filename=path.name)

    try:
        from markitdown import MarkItDown

        result = MarkItDown().convert(str(path))
        if result.text_content and result.text_content.strip():
            return result.text_content
    except Exception as exc:
        logger.warning("parse_docx_failed", filename=path.name, error=str(exc))

    raise FileParseError(path.name, reason="DOCX 解析失败。")


async def parse_image(file_path: str | Path) -> str:
    """把图片解析为 Markdown。"""

    from app.core.llm import acompletion

    path = Path(file_path)
    logger.info("parse_image_start", filename=path.name)

    mime_map = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".webp": "image/webp",
    }
    mime_type = mime_map.get(path.suffix.lower(), "image/png")
    encoded = base64.b64encode(path.read_bytes()).decode("utf-8")
    prompt = populate_prompt(SYSTEM_PROMPT_IMAGE_PARSE)
    messages: list[ChatMessage] = [
        {
            "role": USER,
            "content": [
                {"type": "text", "text": prompt},
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:{mime_type};base64,{encoded}"},
                },
            ],
        }
    ]

    try:
        text = await acompletion(messages)
    except Exception as exc:
        logger.error("parse_image_failed", filename=path.name, error=str(exc))
        raise FileParseError(path.name, reason=f"图片解析失败：{exc}") from exc

    if not text.strip():
        raise FileParseError(path.name, reason="图片解析结果为空。")
    return text
