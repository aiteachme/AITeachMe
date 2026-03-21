"""PPTX 解析器：文本 + 图片提取。"""

from __future__ import annotations

from pathlib import Path

import structlog

from app.agents.ingest.parse_utils import save_image_bytes
from app.core.exceptions import FileParseError

logger = structlog.get_logger()


async def parse_pptx(file_path: str | Path, asset_dir: Path) -> str:
    """PPTX 解析主入口：MarkItDown 优先，python-pptx 兜底。"""
    path = Path(file_path)
    logger.info("parse_pptx_start", filename=path.name)

    # 先尝试 MarkItDown（markdown 结构更好）
    try:
        from markitdown import MarkItDown

        result = MarkItDown().convert(str(path))
        if result.text_content and result.text_content.strip():
            _extract_pptx_images(path, asset_dir)
            return result.text_content
    except Exception as exc:
        logger.warning("parse_pptx_markitdown_failed", filename=path.name, error=str(exc))

    # Fallback: python-pptx 手动提取
    try:
        return _parse_pptx_native(path, asset_dir)
    except Exception as exc:
        logger.warning("parse_pptx_native_failed", filename=path.name, error=str(exc))

    raise FileParseError(path.name, reason="PPTX 解析失败。")


# ---------------------------------------------------------------------------
# python-pptx 原生解析
# ---------------------------------------------------------------------------


def _parse_pptx_native(file_path: Path, asset_dir: Path) -> str:
    """用 python-pptx 手动提取 slide 文本和图片。"""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(file_path))
    sections: list[str] = []
    image_count = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        sections.append(f"## 幻灯片 {slide_idx}")

        for shape in slide.shapes:
            # 文本
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        sections.append(text)

            # 图片
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count = _try_extract_shape_image(
                    shape, asset_dir, slide_idx, image_count, sections,
                )

            # 组合图形中的图片
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_count = _try_extract_shape_image(
                            sub_shape, asset_dir, slide_idx, image_count, sections,
                        )

        sections.append("")

    result = "\n".join(sections).strip()
    if not result:
        raise FileParseError(file_path.name, reason="PPTX 文本提取为空。")

    logger.info("parse_pptx_native_done", filename=file_path.name, images=image_count)
    return result


def _try_extract_shape_image(
    shape: object, asset_dir: Path, slide_idx: int, image_count: int, sections: list[str],
) -> int:
    """尝试从 shape 中提取图片，成功则追加 markdown 引用。返回更新后的 image_count。"""
    try:
        image = shape.image  # type: ignore[attr-defined]
        img_bytes = image.blob
        ext = image.content_type.split("/")[-1] if image.content_type else "png"
        if ext == "jpeg":
            ext = "jpg"
        if len(img_bytes) < 1024:
            return image_count
        image_count += 1
        filename = save_image_bytes(
            img_bytes, asset_dir,
            name_hint=f"slide{slide_idx}_img{image_count}",
            ext=f".{ext}",
        )
        sections.append(f"\n![幻灯片{slide_idx} 图片{image_count}]({filename})\n")
    except Exception:
        logger.debug("pptx_image_extract_failed", slide=slide_idx)
    return image_count


# ---------------------------------------------------------------------------
# 图片补充提取
# ---------------------------------------------------------------------------


def _extract_pptx_images(file_path: Path, asset_dir: Path) -> None:
    """从 PPTX 中补充提取图片（当 MarkItDown 成功时调用）。"""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return

    prs = Presentation(str(file_path))
    count = 0
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob
                    ext = shape.image.content_type.split("/")[-1] if shape.image.content_type else "png"
                    if ext == "jpeg":
                        ext = "jpg"
                    if len(img_bytes) < 1024:
                        continue
                    save_image_bytes(
                        img_bytes, asset_dir,
                        name_hint=f"slide{slide_idx}_img{count}",
                        ext=f".{ext}",
                    )
                    count += 1
                except Exception:
                    pass
    if count > 0:
        logger.info("pptx_images_supplemented", filename=file_path.name, count=count)
